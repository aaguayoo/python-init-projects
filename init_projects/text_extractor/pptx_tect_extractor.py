"""PPTX Text Extractor
This file contains the function needed for extracting text from a PPTX file."""

from pptx import Presentation


def extract_text_ppptx(file: str) -> str:
    """
    This function extracts the text from a PPTX file.
    Args:
        file (str): PPTX file.
    Returns:
        str: The extracted text from the PPTX file.
    Raises:
        ValueError: If the file is not a PPTX.
    """
    try:
        file = Presentation(file)
    except Exception:
        raise ValueError("The file is not a valid PPTX.")

    text = ""
    for slide in file.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
    return text


extract_text_ppptx("texto_prueba_pptx.pptx")
