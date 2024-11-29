"""PPTX Text Extractor
This file contains the function needed for extracting text from a PPTX file."""

from pptx import Presentation


def extract_text_ppptx(file: str) -> str:

    """This funtion extracts the text from PPPTX file.

    The ecope of this funtion is extract text from PPPTX file.

    Args:
        arg_1(str):
             DOCX file.

    Returns:
        str: The extracted text from PPPTX file.

    Raises:
        ValueError: If fike is not a PPPTX file.
    """

    try:
        file = Presentation(file)
    except Exception:
        raise ValueError("The file is not a valid PPTX.")

    text = ""
    for slide in file.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
    return text
